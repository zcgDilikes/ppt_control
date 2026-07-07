"""PPT command executor — thin ``pyautogui`` / ``pyperclip`` wrapper.

Routes incoming command dicts (already decoded by ``ppt_core.ws_messages``
and routed via ``ppt_core.command_dispatcher.CommandDispatcher``) to the
appropriate OS-level keyboard / clipboard / file open action.

Supported commands (also enumerated in
``ppt_core.command_dispatcher.CommandDispatcher.PPT_CMDS``):

* ``NEXT_PAGE``         — ``pagedown``
* ``PREV_PAGE``         — ``pageup``
* ``FULL_SCREEN``       — ``f5``
* ``FROM_CURRENT``      — ``shift + f5``
* ``BLACK_SCREEN``      — ``b``
* ``WHITE_SCREEN``      — ``w``
* ``EXIT``              — ``esc``
* ``SEND_TEXT``         — copies ``text`` to the clipboard then ``Ctrl+V``
* ``SELECT_ALL``        — ``Ctrl+A``
* ``COPY``              — ``Ctrl+C``
* ``PASTE``             — ``Ctrl+V``
* ``DELETE``            — ``backspace``
* ``SCREENSHOT``        — saves a screenshot to ``save_dir/screen_<ts>.png``
* ``OPEN_PPT``          — opens an existing ``.pptx``/``...`` file or creates
                          a new empty ``.pptx`` in the system temp dir

All ``pyautogui`` / ``pyperclip`` imports are wrapped in try/except so the
module imports cleanly even when those packages are not installed — every
action that requires them is also guarded with try/except so the executor
fails silently rather than raising into the dispatcher's lock.
"""

from __future__ import annotations

import os
import tempfile
import time
from typing import Callable, Optional

# ``pyautogui`` is the keyboard / mouse / screenshot driver.  It is the
# heaviest third-party dependency in the stack; we import it defensively so
# that a missing installation doesn't prevent the executor module from
# being imported at all (e.g. during ``python -c "import ..."`` smoke tests
# on a CI box without GUI libraries).
try:
    import pyautogui as _pyautogui  # type: ignore[import-untyped]
except Exception:  # pragma: no cover - exercised only without pyautogui
    _pyautogui = None  # type: ignore[assignment]

# ``pyperclip`` is used to push text into the clipboard for ``SEND_TEXT``.
# It is optional for the rest of the executor — only ``SEND_TEXT`` needs
# it, and that path is wrapped in try/except.
try:
    import pyperclip as _pyperclip  # type: ignore[import-untyped]
except Exception:  # pragma: no cover - exercised only without pyperclip
    _pyperclip = None  # type: ignore[assignment]

# ``win32com`` + ``pythoncom`` are used to drive PowerPoint directly via its
# COM API. This bypasses the focus issue that plagues pyautogui when our
# Qt app has keyboard focus — synthetic keys sent via pyautogui go to the
# focused window (Qt), not the slideshow window (PowerPoint). COM drives
# the slideshow regardless of focus. Patterned after ppt_notes._ensure_pywin32.
_wc = None  # type: ignore[var-annotated]


def _ensure_pywin32() -> bool:
    """Lazily import ``win32com.client``. Returns True on success."""
    global _wc
    if _wc is not None:
        return True
    try:
        import win32com.client as wc  # type: ignore[import-not-found]
    except ImportError:
        return False
    _wc = wc
    return True


def _ppt_show_view():
    """Return the active PowerPoint slideshow View, or None if not available.

    Walks ``GetObject("PowerPoint.Application")`` → ``SlideShowWindows`` →
    first window's ``View``. Returns None on any failure (PPT not running,
    no active slideshow, pywin32 missing).
    """
    if not _ensure_pywin32():
        return None
    try:
        # error.txt [7]:用 GetActiveObject 替代 GetObject(class=),只获取已运行的
        # PPT 实例,不会创建幽灵 PowerPoint.exe 进程。
        # GetObject(path, clsctx) 会在没运行时创建空实例。
        clsctx = getattr(_wc, "CLSCTX_LOCAL_SERVER", 0x4)
        app = _wc.GetActiveObject("PowerPoint.Application", clsctx)
        windows = app.SlideShowWindows
        if windows.Count < 1:
            return None
        return windows.Item(1).View
    except Exception:
        return None


# SlideShowState enum (msotn.OL 12.0). Hard-coded to avoid the COM typelib
# import requirement at module load.
_PP_STATE_BLACK_SCREEN = 9   # ppSlideShowBlackScreen
_PP_STATE_WHITE_SCREEN = 10  # ppSlideShowWhiteScreen


# PowerPoint file extensions accepted by ``OPEN_PPT``.  Provided as a module
# constant so callers (e.g. file upload validators) can import it without
# instantiating the executor.
PPT_EXTS = {".ppt", ".pptx", ".pptm", ".pps", ".ppsx", ".pot", ".potx"}


class PptExecutor:
    """Execute PPT / slideshow commands via ``pyautogui`` and ``pyperclip``.

    The executor holds no per-instance state beyond ``save_dir`` and the
    optional ``on_screenshot`` callback.  Each ``execute`` call is fire-and-
    forget — there is no queue, no async dispatch, and no return value.

    Parameters
    ----------
    save_dir
        Directory in which to write screenshots produced by the
        ``SCREENSHOT`` command.  Created lazily on first screenshot.
    on_screenshot
        Optional callable invoked with the absolute path of each saved
        screenshot after a successful save.  Failures inside the callback
        are swallowed so that executor never raises into the dispatcher.
    """

    def __init__(
        self,
        *,
        save_dir: str = "./ppt_files/",
        on_screenshot: Optional[Callable[[str], None]] = None,
    ) -> None:
        self._save_dir = save_dir
        self._on_screenshot = on_screenshot

    # ------------------------------------------------------------------ public

    def execute(self, d: dict) -> None:  # noqa: C901 - dispatch table is clearer linear
        """Dispatch one command dict to the matching OS-level action.

        Unknown commands are silently ignored.  Every action is guarded
        against ``pyautogui`` / ``pyperclip`` being unavailable (those
        imports may have failed at module load) and against the underlying
        call itself raising; the executor must never propagate exceptions
        into ``CommandDispatcher.dispatch``'s lock.
        """
        cmd = d.get("cmd")

        if cmd == "NEXT_PAGE":
            # COM 优先:不受焦点问题影响;fallback 到 pyautogui("pagedown")
            if not self._show_view_next():
                self._press("pagedown")
            return

        if cmd == "PREV_PAGE":
            if not self._show_view_previous():
                self._press("pageup")
            return

        if cmd == "FULL_SCREEN":
            # 没有活跃放映 → 启动;已在放映 → 退出(避免重复触发)
            if _ppt_show_view() is not None:
                if not self._show_view_exit():
                    self._press("esc")
            else:
                self._press("f5")
            return

        if cmd == "FROM_CURRENT":
            self._hotkey("shift", "f5")
            return

        if cmd == "BLACK_SCREEN":
            if not self._show_view_set_state(_PP_STATE_BLACK_SCREEN):
                self._press("b")
            return

        if cmd == "WHITE_SCREEN":
            if not self._show_view_set_state(_PP_STATE_WHITE_SCREEN):
                self._press("w")
            return

        if cmd == "EXIT":
            if not self._show_view_exit():
                self._press("esc")
            return

        if cmd == "SEND_TEXT":
            self._send_text(d.get("text", ""))
            return

        if cmd == "SELECT_ALL":
            self._hotkey("ctrl", "a")
            return

        if cmd == "COPY":
            self._hotkey("ctrl", "c")
            return

        if cmd == "PASTE":
            self._hotkey("ctrl", "v")
            return

        if cmd == "DELETE":
            # The brief calls for try/except around DELETE explicitly.
            try:
                self._press("backspace")
            except Exception:
                pass
            return

        # info.txt 一.2:MOUSE_DOWN/UP 支持捏合长按拖拽
        if cmd == "MOUSE_DOWN":
            self._mouse_down()
            return
        if cmd == "MOUSE_UP":
            self._mouse_up()
            return

        if cmd == "SCREENSHOT":
            self._screenshot()
            return

        if cmd == "OPEN_PPT":
            self._open_ppt(d.get("path", ""))
            return

        # Unknown cmd — ignored silently, per dispatcher contract.

    # ----------------------------------------------------------------- private

    def _show_view_next(self) -> bool:
        """COM ``SlideShowWindows(1).View.Next()``. True on success."""
        view = _ppt_show_view()
        if view is None:
            return False
        try:
            view.Next()
            return True
        except Exception:
            return False

    def _show_view_previous(self) -> bool:
        view = _ppt_show_view()
        if view is None:
            return False
        try:
            view.Previous()
            return True
        except Exception:
            return False

    def _show_view_exit(self) -> bool:
        view = _ppt_show_view()
        if view is None:
            return False
        try:
            view.Exit()
            return True
        except Exception:
            return False

    def _show_view_set_state(self, state: int) -> bool:
        """COM ``View.State = <state>``. Used for black/white screen."""
        view = _ppt_show_view()
        if view is None:
            return False
        try:
            view.State = state
            return True
        except Exception:
            return False

    def _mouse_down(self) -> None:
        """Hold the left mouse button down (for drag/select)."""
        if _pyautogui is None:
            return
        try:
            _pyautogui.mouseDown()
        except Exception:
            pass

    def _mouse_up(self) -> None:
        """Release the left mouse button (for drag/select end)."""
        if _pyautogui is None:
            return
        try:
            _pyautogui.mouseUp()
        except Exception:
            pass

    def _press(self, key: str) -> None:
        if _pyautogui is None:
            return
        try:
            _pyautogui.press(key)
        except Exception:
            pass

    def _hotkey(self, *keys: str) -> None:
        if _pyautogui is None:
            return
        try:
            _pyautogui.hotkey(*keys)
        except Exception:
            pass

    def _send_text(self, text: str) -> None:
        """Push ``text`` into the clipboard then trigger ``Ctrl+V``.

        Both the clipboard write and the paste hotkey are guarded so the
        executor still works (silently doing nothing) when ``pyperclip``
        or ``pyautogui`` are unavailable.
        """
        if not isinstance(text, str):
            try:
                text = str(text)
            except Exception:
                return
        if _pyperclip is not None:
            try:
                _pyperclip.copy(text)
            except Exception:
                # No clipboard available — fall through and try the paste
                # anyway; on most systems that is a no-op rather than an
                # error, and we have nothing else to do here.
                pass
        self._hotkey("ctrl", "v")

    def _screenshot(self) -> None:
        """Capture the screen, save it, and notify ``on_screenshot``.

        The filename embeds ``int(time.time())`` so successive screenshots
        in the same second collide — that matches the brief and is fine
        for the typical "one screenshot per gesture" workload.  The
        directory is created lazily; ``on_screenshot`` failures are
        swallowed so the executor remains total.
        """
        if _pyautogui is None:
            return
        out_path = os.path.join(self._save_dir, f"screen_{int(time.time())}.png")
        abs_path: Optional[str] = None
        try:
            try:
                os.makedirs(self._save_dir, exist_ok=True)
            except Exception:
                # If we cannot create the directory, bail — there is no
                # useful recovery for a screenshot that cannot be written.
                return
            try:
                _pyautogui.screenshot(out_path)
            except Exception:
                return
            abs_path = os.path.abspath(out_path)
        except Exception:
            return
        if abs_path is not None and self._on_screenshot is not None:
            try:
                self._on_screenshot(abs_path)
            except Exception:
                pass

    def _open_ppt(self, path: str) -> None:
        """Open ``path`` if it exists, else launch a fresh empty ``.pptx``.

        Resolution order:

        1. If ``path`` is a non-empty string that points to an existing
           file with a known PPT extension (or any extension, since the
           caller already validated it), ``os.startfile`` is invoked.
        2. Otherwise create a real empty ``.pptx`` (valid OOXML/ZIP
           archive containing a single blank slide) in the system temp
           dir and open that.
        3. ``os.startfile`` is Windows-only; on other platforms this
           silently does nothing because the project as a whole is
           Windows-only (``ppt_pc_client`` is a desktop app for Windows).

        error.txt [1]: 之前用 mkstemp 创建 0 字节 .pptx 空文件,PowerPoint 打开
        失败,被 except 吞掉,用户无感知。改用 python-pptx 创建有效空白演示文稿。
        """
        target: Optional[str] = None

        if isinstance(path, str) and path:
            try:
                if os.path.isfile(path):
                    target = path
            except Exception:
                target = None

        if target is None:
            try:
                target = self._create_blank_pptx()
            except Exception:
                return

        if target is None:
            return

        try:
            os.startfile(target)  # type: ignore[attr-defined]  # Windows-only
        except Exception:
            # ``os.startfile`` may fail for a number of reasons (no file
            # association, no PowerPoint installed, etc).  Swallow —
            # caller cannot meaningfully react, and the executor's
            # contract is total.
            pass

    @staticmethod
    def _create_blank_pptx() -> Optional[str]:
        """Create a real empty .pptx file in the temp dir.

        Tries python-pptx first; falls back to COM
        (PowerPoint.Application.Presentations.Add) if available;
        falls back to copying a minimal OOXML template (or returning
        None if nothing works — caller will silently skip).

        error.txt [1]: 0 字节 .pptx 在 PowerPoint 打开会失败。这里必须返回
        合法 OOXML 文件才能被 os.startfile 正确处理。
        """
        temp_root = os.environ.get("TEMP") or tempfile.gettempdir()
        try:
            fd, target = tempfile.mkstemp(suffix=".pptx", prefix="ppt_remote_", dir=temp_root)
            try:
                os.close(fd)
            except Exception:
                pass
        except Exception:
            return None

        # Strategy 1: python-pptx
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation()
            # Add one blank slide so the file is non-empty and valid.
            try:
                blank_layout = prs.slide_layouts[6]
            except (IndexError, AttributeError):
                blank_layout = None
            if blank_layout is not None:
                prs.slides.add_slide(blank_layout)
            prs.save(target)
            return target
        except Exception:
            pass  # fall through to COM

        # Strategy 2: COM (PowerPoint must be installed)
        try:
            import win32com.client  # type: ignore
            app = win32com.client.DispatchEx("PowerPoint.Application")
            try:
                app.Visible = True
                pres = app.Presentations.Add()
                try:
                    pres.SaveAs(target)
                finally:
                    pres.Close()
            finally:
                try:
                    app.Quit()
                except Exception:
                    pass
            return target
        except Exception:
            pass

        # Strategy 3: all failed. Remove the empty file and return None.
        try:
            os.unlink(target)
        except Exception:
            pass
        return None
